"""
ETL Phase 1: 文档解析 + 页面渲染

核心能力：
1. extract_slides() — 使用 python-pptx 提取每页文本（标题/正文/备注）
2. render_slides_to_images() — LibreOffice → PDF → PNG 渲染缩略图
3. compute_checksum() — 文件级 MD5 去重
"""
from __future__ import annotations

import hashlib
import logging
import os
import subprocess
import tempfile
import uuid
from pathlib import Path
from typing import Any

from pptx import Presentation

logger = logging.getLogger(__name__)


# ═══════════════════════════════════════════════════════════
# 1. 文本提取
# ═══════════════════════════════════════════════════════════


def extract_slides(pptx_path: str | Path) -> list[dict[str, Any]]:
    """
    使用 python-pptx 提取 PPT 所有页面的文本内容。

    返回值：list[dict]，每个 dict 包含：
        - slide_index: int（从 1 开始）
        - title: str（优先取幻灯片标题占位符，其次取第一个段落）
        - body_text: str（非标题的正文文本）
        - notes_text: str（页面备注）
        - text_length: int（纯文本总长度，用于后续判断是否纯图片页）
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        logger.error(f"文件不存在: {pptx_path}")
        return []

    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        logger.error(f"无法打开 PPT 文件: {pptx_path}, 错误: {e}")
        return []

    slides_data: list[dict[str, Any]] = []

    for idx, slide in enumerate(prs.slides, start=1):
        title_text = ""
        body_parts: list[str] = []
        notes_text = ""

        # 提取标题和正文
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                # 判断是否为标题占位符
                if shape.is_placeholder and shape.placeholder_format.idx == 0:
                    title_text = text
                elif shape.name.startswith("Title"):
                    title_text = text
                else:
                    body_parts.append(text)

        # 如果没找到标题，取第一个非空段落作为标题
        if not title_text and body_parts:
            title_text = body_parts[0]
            body_parts = body_parts[1:]

        # 提取备注
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text.strip() if notes_slide.notes_text_frame else ""

        body_text = "\n".join(body_parts)
        text_length = len(title_text) + len(body_text) + len(notes_text)

        slides_data.append({
            "slide_index": idx,
            "title": title_text,
            "body_text": body_text,
            "notes_text": notes_text,
            "text_length": text_length,
        })

    logger.info(f"提取完成: {pptx_path.name} → {len(slides_data)} 页")
    return slides_data


# ═══════════════════════════════════════════════════════════
# 2. 页面渲染（PPT → PDF → PNG）
# ═══════════════════════════════════════════════════════════


def render_slides_to_images(
    pptx_path: str | Path,
    output_dir: str | Path,
    dpi: int = 150,
) -> list[str]:
    """
    将 PPT 每页渲染为 PNG 缩略图。

    流程：PPT → PDF（LibreOffice 命令行）→ PNG（pdf2image）

    参数：
        pptx_path: 源 PPT 文件路径
        output_dir: 输出目录（每页保存为 slide_001.png, slide_002.png ...）
        dpi: 渲染分辨率（默认 150，越大质量越高但文件越大）

    返回：
        list[str] — 生成的 PNG 文件路径列表（按页码顺序）

    异常：
        RuntimeError: LibreOffice 或 pdf2image 不可用
    """
    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    # ── Step 1: PPT → PDF ──────────────────────────────────
    with tempfile.TemporaryDirectory() as tmpdir:
        pdf_path = os.path.join(tmpdir, f"{pptx_path.stem}.pdf")

        try:
            result = subprocess.run(
                [
                    "libreoffice", "--headless", "--convert-to", "pdf",
                    "--outdir", tmpdir,
                    str(pptx_path.resolve()),
                ],
                capture_output=True,
                text=True,
                timeout=120,
            )
            if result.returncode != 0:
                raise RuntimeError(
                    f"LibreOffice 转换失败 (exit={result.returncode}): "
                    f"{result.stderr.strip()}"
                )
            # LibreOffice 输出 PDF 文件名可能与源文件同名
            expected_pdf = os.path.join(tmpdir, f"{pptx_path.stem}.pdf")
            if not os.path.exists(expected_pdf):
                # 尝试查找生成的 pdf
                pdf_files = [f for f in os.listdir(tmpdir) if f.endswith(".pdf")]
                if pdf_files:
                    expected_pdf = os.path.join(tmpdir, pdf_files[0])
                else:
                    raise RuntimeError("LibreOffice 未生成 PDF 文件")
            pdf_path = expected_pdf
        except FileNotFoundError:
            raise RuntimeError("LibreOffice 未安装，请执行: sudo apt install libreoffice-core libreoffice-impress")
        except subprocess.TimeoutExpired:
            raise RuntimeError("LibreOffice 转换超时（120s），文件可能过大或损坏")

        # ── Step 2: PDF → PNG ─────────────────────────────
        try:
            from pdf2image import convert_from_path
        except ImportError:
            raise RuntimeError("pdf2image 未安装，请执行: pip install pdf2image")

        try:
            images = convert_from_path(pdf_path, dpi=dpi)
        except Exception as e:
            raise RuntimeError(f"PDF 转 PNG 失败: {e}")

        # ── Step 3: 保存 PNG ──────────────────────────────
        png_paths: list[str] = []
        for i, img in enumerate(images, start=1):
            png_name = f"slide_{i:03d}.png"
            png_path = str(output_dir / png_name)
            img.save(png_path, "PNG")
            png_paths.append(png_path)

        logger.info(f"渲染完成: {pptx_path.name} → {len(png_paths)} 张缩略图 → {output_dir}")
        return png_paths


# ═══════════════════════════════════════════════════════════
# 3. 文件级去重
# ═══════════════════════════════════════════════════════════


def compute_checksum(file_path: str | Path) -> str:
    """
    计算文件 MD5 checksum（用于去重）。

    以 64KB 块为单位流式读取，避免大文件内存占用过高。
    """
    file_path = Path(file_path)
    if not file_path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")

    md5 = hashlib.md5()
    with open(file_path, "rb") as f:
        while chunk := f.read(65536):  # 64KB 块
            md5.update(chunk)
    return md5.hexdigest()


# ═══════════════════════════════════════════════════════════
# 4. 页面级去重（simhash 近似判断）
# ═══════════════════════════════════════════════════════════


def _tokenize(text: str) -> list[str]:
    """简单中文分词：按标点和空格分割"""
    import re
    tokens = re.split(r'[\s,，。；;：:！!？?、()（）【】\[\]{}""''\n\r]+', text)
    return [t.strip() for t in tokens if t.strip()]


def simhash(text: str, hash_bits: int = 64) -> int:
    """计算文本的 simhash 值"""
    tokens = _tokenize(text)
    if not tokens:
        return 0

    v = [0] * hash_bits
    for token in tokens:
        # 使用 Python 内置 hash（确定性）
        h = hash(token) & ((1 << hash_bits) - 1)
        for i in range(hash_bits):
            if h & (1 << i):
                v[i] += 1
            else:
                v[i] -= 1

    fingerprint = 0
    for i in range(hash_bits):
        if v[i] > 0:
            fingerprint |= (1 << i)
    return fingerprint


def simhash_similarity(h1: int, h2: int, hash_bits: int = 64) -> float:
    """计算两个 simhash 的海明距离相似度（0~1）"""
    if h1 == 0 and h2 == 0:
        return 1.0
    xor_val = h1 ^ h2
    distance = bin(xor_val).count("1")
    return 1.0 - (distance / hash_bits)


def check_slide_dedup(
    body_text: str,
    existing_simhashes: list[int],
    threshold: float = 0.9,
) -> bool:
    """
    检查页面是否与已有页面近似重复。

    参数：
        body_text: 当前页面的正文文本
        existing_simhashes: 已有页面的 simhash 列表
        threshold: 相似度阈值（默认 0.9，即海明距离 ≤ 6）

    返回：
        True — 近似重复，应跳过
        False — 不重复，可以入库
    """
    if not body_text.strip():
        return False  # 纯图片页不参与去重
    h = simhash(body_text)
    for existing in existing_simhashes:
        if simhash_similarity(h, existing) >= threshold:
            return True
    return False
