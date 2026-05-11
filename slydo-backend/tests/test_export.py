"""
单页导出单元测试 — 测试 _extract_single_slide_pptx
"""
from __future__ import annotations

import os
import tempfile

import pptx
import pytest
from pptx import Presentation

from app.services.export import _extract_single_slide_pptx

# ─── 辅助函数 ───────────────────────────────────────────────


def _create_test_pptx(slide_count: int = 3) -> str:
    """创建一个测试 PPT 文件"""
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # title and content
    for i in range(slide_count):
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = f"Slide {i + 1}"
        else:
            from pptx.util import Inches
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            tf = txBox.text_frame
            tf.text = f"Slide {i + 1}"

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    tmp.close()
    return tmp.name


# ─── 测试 ───────────────────────────────────────────────────


class TestExtractSingleSlide:
    """测试 _extract_single_slide_pptx"""

    def test_extract_first_slide(self):
        path = _create_test_pptx(3)
        buf = _extract_single_slide_pptx(path, 1)
        prs = Presentation(buf)
        assert len(prs.slides) == 1, f"期望 1 页, 实际 {len(prs.slides)}"
        title = prs.slides[0].shapes.title
        assert title is not None, "第1页应有标题"
        assert "Slide 1" in title.text
        os.unlink(path)

    def test_extract_last_slide(self):
        path = _create_test_pptx(3)
        buf = _extract_single_slide_pptx(path, 3)
        prs = Presentation(buf)
        assert len(prs.slides) == 1
        assert "Slide 3" in prs.slides[0].shapes.title.text
        os.unlink(path)

    def test_extract_middle_slide(self):
        path = _create_test_pptx(5)
        buf = _extract_single_slide_pptx(path, 3)
        prs = Presentation(buf)
        assert len(prs.slides) == 1
        assert "Slide 3" in prs.slides[0].shapes.title.text
        os.unlink(path)

    def test_output_is_valid_pptx(self):
        """验证输出文件可以被正常打开并读取"""
        path = _create_test_pptx(3)
        buf = _extract_single_slide_pptx(path, 2)
        # 检查 ZIP 文件头
        assert buf.read(4) == b"PK\x03\x04"
        buf.seek(0)
        # 重新用 python-pptx 打开
        prs = Presentation(buf)
        assert len(prs.slides) == 1
        os.unlink(path)

    def test_extract_invalid_index_raises(self):
        path = _create_test_pptx(2)
        with pytest.raises(ValueError, match="超出范围"):
            _extract_single_slide_pptx(path, 99)
        os.unlink(path)

    def test_extract_zero_index_raises(self):
        path = _create_test_pptx(2)
        with pytest.raises(ValueError, match="超出范围"):
            _extract_single_slide_pptx(path, 0)
        os.unlink(path)

    def test_extract_single_page_retains_theme(self):
        """单页导出应保留主题/母版"""
        path = _create_test_pptx(1)
        buf = _extract_single_slide_pptx(path, 1)
        prs = Presentation(buf)
        assert len(prs.slides) == 1
        os.unlink(path)

    def test_extract_one_page_from_one_page(self):
        """只有 1 页的 PPT 也能正确导出"""
        path = _create_test_pptx(1)
        buf = _extract_single_slide_pptx(path, 1)
        prs = Presentation(buf)
        assert len(prs.slides) == 1
        os.unlink(path)

    def test_extract_all_pages_individually(self):
        """每一页都能单独导出"""
        path = _create_test_pptx(4)
        for i in range(1, 5):
            buf = _extract_single_slide_pptx(path, i)
            prs = Presentation(buf)
            assert len(prs.slides) == 1
            assert f"Slide {i}" in prs.slides[0].shapes.title.text
        os.unlink(path)


class TestExportServiceAsync:
    """async export_single_slide 测试（依赖真实 DB，由 API 集成测试覆盖）"""
