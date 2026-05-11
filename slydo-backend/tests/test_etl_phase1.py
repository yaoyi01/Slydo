"""
ETL Phase 1 单元测试 —— 文档解析 + 渲染 + 去重
"""
from __future__ import annotations

import hashlib
import os
import tempfile
from pathlib import Path

import pytest

from app.services.etl.phase1_extract import (
    compute_checksum,
    extract_slides,
    render_slides_to_images,
    simhash,
    simhash_similarity,
)

# ── 测试用 PPT 文件路径 ──────────────────────────────────
TEST_PPTX = Path(
    "/mnt/c/Users/kiven/Documents/LLM_wiki/知识库/raw/sources/信息防扩散解决方案.pptx"
)
HAS_REAL_PPTX = TEST_PPTX.exists()


# ═══════════════════════════════════════════════════════════
# extract_slides() 测试
# ═══════════════════════════════════════════════════════════


class TestExtractSlides:
    """文档解析测试"""

    def test_normal_pptx_extracts_all_slides(self):
        """正常 PPT 能正确提取所有页面"""
        if not HAS_REAL_PPTX:
            pytest.skip("测试 PPT 文件不存在")
        result = extract_slides(TEST_PPTX)
        assert len(result) > 0, "应提取到至少 1 页"
        assert len(result) == 36, "该文件应有 36 页"

    def test_slide_has_required_fields(self):
        """每页数据结构完整"""
        if not HAS_REAL_PPTX:
            pytest.skip("测试 PPT 文件不存在")
        slides = extract_slides(TEST_PPTX)
        for s in slides:
            assert "slide_index" in s
            assert "title" in s
            assert "body_text" in s
            assert "notes_text" in s
            assert "text_length" in s
            assert isinstance(s["slide_index"], int)
            assert s["slide_index"] >= 1

    def test_first_slide_has_title(self):
        """第一页应有标题"""
        if not HAS_REAL_PPTX:
            pytest.skip("测试 PPT 文件不存在")
        slides = extract_slides(TEST_PPTX)
        assert slides[0]["title"] != "", "第一页应有标题"

    def test_slide_with_notes_extracted(self):
        """有备注的页面应提取到备注文本"""
        if not HAS_REAL_PPTX:
            pytest.skip("测试 PPT 文件不存在")
        slides = extract_slides(TEST_PPTX)
        slides_with_notes = [s for s in slides if s["notes_text"]]
        assert len(slides_with_notes) > 0, "应至少有一页包含备注"

    def test_corrupted_file_returns_empty_list(self):
        """损坏文件返回空列表而非抛异常"""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            f.write(b"not a valid pptx file")
            tmp_path = f.name
        try:
            result = extract_slides(tmp_path)
            assert result == []
        finally:
            os.unlink(tmp_path)

    def test_nonexistent_file_returns_empty_list(self):
        """不存在的文件返回空列表"""
        result = extract_slides("/tmp/nonexistent_file.pptx")
        assert result == []

    def test_image_only_slide_returns_empty_title(self):
        """纯图片页提取到空标题和空正文"""
        from app.services.etl.phase1_extract import extract_slides
        test_file = "/tmp/test_image_only.pptx"
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image
        import io
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img = Image.new('RGB', (10, 10), color='red')
        buf = io.BytesIO()
        img.save(buf, format='PNG')
        slide.shapes.add_picture(buf, Inches(0), Inches(0), Inches(1), Inches(1))
        prs.save(test_file)
        try:
            slides = extract_slides(test_file)
            assert len(slides) == 1, "纯图片页应返回 1 页"
            s = slides[0]
            assert s["title"] == "", "纯图片页 title 应为空"
            assert s["body_text"] == "", "纯图片页 body_text 应为空"
            assert s["text_length"] == 0, "纯图片页 text_length 应为 0"
        finally:
            import os
            os.unlink(test_file)


# ═══════════════════════════════════════════════════════════
# render_slides_to_images() 测试
# ═══════════════════════════════════════════════════════════


class TestRenderSlidesToImages:
    """页面渲染测试"""

    def test_renders_all_pages(self):
        """正常 PPT 每页都生成对应的 PNG 缩略图"""
        if not HAS_REAL_PPTX:
            pytest.skip("测试 PPT 文件不存在")
        with tempfile.TemporaryDirectory() as tmpdir:
            pngs = render_slides_to_images(TEST_PPTX, tmpdir, dpi=100)
            # 检查数量
            assert len(pngs) == 36, f"应生成 36 张缩略图, 实际 {len(pngs)}"
            # 检查文件命名
            assert os.path.basename(pngs[0]) == "slide_001.png"
            assert os.path.basename(pngs[-1]) == "slide_036.png"
            # 检查文件大小（不应为空）
            for p in pngs:
                assert os.path.getsize(p) > 1000, f"缩略图太小: {p}"

    def test_render_output_dir_is_created(self):
        """输出目录自动创建"""
        if not HAS_REAL_PPTX:
            pytest.skip("测试 PPT 文件不存在")
        with tempfile.TemporaryDirectory() as tmpdir:
            subdir = os.path.join(tmpdir, "nested", "dir")
            pngs = render_slides_to_images(TEST_PPTX, subdir, dpi=100)
            assert os.path.isdir(subdir)
            assert len(pngs) > 0

    def test_render_low_dpi_smaller_files(self):
        """低 DPI 产生的文件更小"""
        if not HAS_REAL_PPTX:
            pytest.skip("测试 PPT 文件不存在")
        with tempfile.TemporaryDirectory() as tmpdir:
            pngs_100 = render_slides_to_images(TEST_PPTX, os.path.join(tmpdir, "dpi100"), dpi=100)
            pngs_200 = render_slides_to_images(TEST_PPTX, os.path.join(tmpdir, "dpi200"), dpi=200)
            size_100 = sum(os.path.getsize(p) for p in pngs_100)
            size_200 = sum(os.path.getsize(p) for p in pngs_200)
            assert size_100 < size_200, "低 DPI 应产生更小的文件"


# ═══════════════════════════════════════════════════════════
# checksum 测试
# ═══════════════════════════════════════════════════════════


class TestComputeChecksum:
    """文件去重：MD5 checksum"""

    def test_deterministic(self):
        """相同文件产生相同的 checksum"""
        if not HAS_REAL_PPTX:
            pytest.skip("测试 PPT 文件不存在")
        cs1 = compute_checksum(TEST_PPTX)
        cs2 = compute_checksum(TEST_PPTX)
        assert cs1 == cs2

    def test_different_files_different_checksums(self):
        """不同文件产生不同的 checksum"""
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as f1:
            f1.write(b"content a")
            p1 = f1.name
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as f2:
            f2.write(b"content b")
            p2 = f2.name
        try:
            cs1 = compute_checksum(p1)
            cs2 = compute_checksum(p2)
            assert cs1 != cs2
        finally:
            os.unlink(p1)
            os.unlink(p2)

    def test_long_file(self):
        """大文件也能正确计算 checksum（流式读取不会 OOM）"""
        with tempfile.NamedTemporaryFile(delete=False, suffix=".bin") as f:
            f.write(b"x" * 10 * 1024 * 1024)  # 10MB
            path = f.name
        try:
            cs = compute_checksum(path)
            expected = hashlib.md5(b"x" * 10 * 1024 * 1024).hexdigest()
            assert cs == expected
        finally:
            os.unlink(path)

    def test_nonexistent_file_raises(self):
        """不存在的文件抛 FileNotFoundError"""
        with pytest.raises(FileNotFoundError):
            compute_checksum("/tmp/nonexistent.pptx")


# ═══════════════════════════════════════════════════════════
# simhash 去重测试
# ═══════════════════════════════════════════════════════════


class TestSimhash:
    """页面级去重：simhash 相似度"""

    def test_identical_texts_are_similar(self):
        """相同文本的 simhash 相似度应为 1.0"""
        text = "联软信息安全解决方案介绍"
        h1 = simhash(text)
        h2 = simhash(text)
        assert simhash_similarity(h1, h2) == 1.0

    def test_slightly_different_texts_still_similar(self):
        """微小差异的文本仍高度相似"""
        base = "联软信息安全解决方案介绍 联软科技 终端安全 数据防泄漏"
        similar = "联软信息安全解决方案介绍 联软科技 终端安全 防泄漏"
        s = simhash_similarity(simhash(base), simhash(similar))
        assert s > 0.7, f"相似文本的相似度应 > 0.7, 实际 {s}"

    def test_completely_different_texts_low_similarity(self):
        """完全不相关的文本相似度很低"""
        t1 = "联软信息安全解决方案介绍"
        t2 = "今天天气真不错出去吃个饭吧"
        s = simhash_similarity(simhash(t1), simhash(t2))
        assert s < 0.65, f"不相关文本的相似度应 < 0.65, 实际 {s}"

    def test_empty_text_returns_zero(self):
        """空文本的 simhash 为 0"""
        assert simhash("") == 0

    def test_check_dedup_returns_true_for_duplicate(self):
        """check_slide_dedup 对近似重复正文返回 True"""
        from app.services.etl.phase1_extract import check_slide_dedup
        body = "终端安全解决方案包括：端点保护、数据防泄漏、准入控制"
        existing = [simhash(body)]
        assert check_slide_dedup(body, existing, threshold=0.9) is True

    def test_check_dedup_returns_false_for_new(self):
        """check_slide_dedup 对不同正文返回 False"""
        from app.services.etl.phase1_extract import check_slide_dedup
        existing = [simhash("终端安全解决方案")]
        assert check_slide_dedup("数据防泄漏产品介绍", existing, threshold=0.9) is False

    def test_check_dedup_empty_text_not_flagged(self):
        """纯图片页（空文本）不被标记为重复"""
        from app.services.etl.phase1_extract import check_slide_dedup
        assert check_slide_dedup("", [], threshold=0.9) is False
